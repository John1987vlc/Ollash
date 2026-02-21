"""
Multi-format document ingester for Knowledge Workspace.
Converts PDF, DOCX, TXT, and Markdown files to indexed text.
"""

from pathlib import Path
from typing import Dict, Optional

from backend.utils.core.system.agent_logger import AgentLogger


class MultiFormatIngester:
    """
    Transforms binary documents (PDF, DOCX) into plain text for indexing.
    Supports: .pdf, .docx, .pptx, .txt, .md, .markdown
    """

    SUPPORTED_FORMATS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}

    def __init__(self, logger: AgentLogger, config: Optional[Dict] = None):
        self.logger = logger
        self.config = config or {}
        self._load_optional_dependencies()

    def _load_optional_dependencies(self):
        """Lazily load PDF/DOCX libraries if available."""
        self.pypdf = None
        self.python_docx = None
        self.python_pptx = None

        try:
            import PyPDF2

            self.pypdf = PyPDF2
            self.logger.debug("PyPDF2 available for PDF processing")
        except ImportError:
            self.logger.warning("PyPDF2 not installed. PDF extraction disabled. Install with: pip install PyPDF2")

        try:
            import docx

            self.python_docx = docx
            self.logger.debug("python-docx available for DOCX processing")
        except ImportError:
            self.logger.warning(
                "python-docx not installed. DOCX extraction disabled. Install with: pip install python-docx"
            )

        try:
            import pptx

            self.python_pptx = pptx
            self.logger.debug("python-pptx available for PPTX processing")
        except ImportError:
            self.logger.debug("python-pptx not installed. PPTX extraction disabled.")

    def ingest_file(self, file_path: Path) -> Optional[str]:
        """
        Ingests a single document file and returns extracted text.
        Returns None if file format is unsupported or extraction fails.
        """
        suffix = file_path.suffix.lower()

        if suffix not in self.SUPPORTED_FORMATS:
            self.logger.warning(f"Unsupported format: {suffix}")
            return None

        try:
            if suffix == ".pdf":
                return self._extract_pdf(file_path)
            elif suffix == ".docx":
                return self._extract_docx(file_path)
            elif suffix == ".pptx":
                return self._extract_pptx(file_path)
            elif suffix in {".txt", ".md", ".markdown"}:
                return self._extract_text(file_path)
        except Exception as e:
            self.logger.error(f"Failed to ingest {file_path.name}: {e}")
            return None

    def _extract_pdf(self, file_path: Path) -> Optional[str]:
        """Extracts text from PDF using PyPDF2."""
        if not self.pypdf:
            self.logger.warning("PyPDF2 not available. Skipping PDF.")
            return None

        text_parts = []
        try:
            with open(file_path, "rb") as f:
                reader = self.pypdf.PdfReader(f)
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{text}")

            return "\n\n".join(text_parts) if text_parts else None
        except Exception as e:
            self.logger.error(f"PDF extraction error ({file_path.name}): {e}")
            return None

    def _extract_docx(self, file_path: Path) -> Optional[str]:
        """Extracts text from DOCX using python-docx."""
        if not self.python_docx:
            self.logger.warning("python-docx not available. Skipping DOCX.")
            return None

        text_parts = []
        try:
            doc = self.python_docx.Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            self.logger.error(f"DOCX extraction error ({file_path.name}): {e}")
            return None

    def _extract_pptx(self, file_path: Path) -> Optional[str]:
        """Extracts text from PPTX using python-pptx."""
        if not self.python_pptx:
            self.logger.debug("python-pptx not available. Skipping PPTX.")
            return None

        text_parts = []
        try:
            prs = self.python_pptx.Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"[Slide {slide_num}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                if slide_text.strip() != f"[Slide {slide_num}]\n":
                    text_parts.append(slide_text)

            return "\n\n".join(text_parts) if text_parts else None
        except Exception as e:
            self.logger.error(f"PPTX extraction error ({file_path.name}): {e}")
            return None

    def _extract_text(self, file_path: Path) -> Optional[str]:
        """Extracts text from plaintext files."""
        try:
            return file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                return file_path.read_text(encoding="latin-1")
            except Exception as e:
                self.logger.error(f"Text extraction error ({file_path.name}): {e}")
                return None

    def ingest_directory(self, directory: Path) -> Dict[str, str]:
        """
        Ingests all supported documents in a directory.
        Returns dict: {file_name: extracted_text}
        """
        results = {}
        for file_path in directory.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_FORMATS:
                text = self.ingest_file(file_path)
                if text:
                    results[file_path.name] = text
                    self.logger.info(f"✓ Ingested: {file_path.name}")

        return results

    def get_file_metadata(self, file_path: Path) -> Dict:
        """Returns metadata about an ingested file."""
        text = self.ingest_file(file_path)
        if not text:
            return {}

        word_count = len(text.split())
        lines = text.count("\n")

        return {
            "file": file_path.name,
            "format": file_path.suffix.lower(),
            "size_bytes": file_path.stat().st_size,
            "word_count": word_count,
            "line_count": lines,
            "extraction_success": True,
        }
